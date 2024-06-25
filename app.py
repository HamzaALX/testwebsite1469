from flask import Flask, render_template, request, send_file
import os
from pdf2docx import Converter
from pptx import Presentation
from pdf2image import convert_from_path
from pytesseract import image_to_string
from PIL import Image
from io import BytesIO
import tempfile
from PyPDF2 import PdfFileMerger
import zipfile
import tabula
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Alignment
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.styles import Font, Alignment
import pytesseract



app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted'
ALLOWED_EXTENSIONS = {'pdf'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/pdf_to_word', methods=['GET', 'POST'])
def pdf_to_word():
    if request.method == 'POST':
        if 'file' not in request.files:
            return render_template('pdf_to_word.html', error='No file part')

        file = request.files['file']
        if file.filename == '':
            return render_template('pdf_to_word.html', error='No selected file')

        if file and allowed_file(file.filename):
            # Save the uploaded file
            filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filename)

            # Convert PDF to Word
            word_filename = os.path.join(app.config['CONVERTED_FOLDER'], file.filename.replace('.pdf', '.docx'))
            cv = Converter(filename)
            cv.convert(word_filename, start=0, end=None)
            cv.close()

            # Return the converted Word file for download
            return send_file(word_filename, as_attachment=True)

    return render_template('pdf_to_word.html')

@app.route('/merge_pdf', methods=['GET', 'POST'])
def merge_pdfs():
    if request.method == 'POST':
        # Get the list of uploaded files
        files = request.files.getlist('file[]')

        if not all(files):
            return render_template('merge_pdf.html', error='Please select files to merge')

        # Save the uploaded files
        filenames = []
        for file in files:
            if file and allowed_file(file.filename):
                filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
                file.save(filename)
                filenames.append(filename)

        if not filenames:
            return render_template('merge_pdf.html', error='No valid PDF files selected')

        # Merge the PDF files
        merged_filename = os.path.join(app.config['CONVERTED_FOLDER'], 'merged_file.pdf')
        merger = PdfFileMerger()
        for filename in filenames:
            merger.append(filename)
        merger.write(merged_filename)
        merger.close()

        # Return the merged PDF file for download
        return send_file(merged_filename, as_attachment=True)

    return render_template('merge_pdf.html')

@app.route('/download/<path:filename>')
def download_file(filename):
    return send_file(filename, as_attachment=True)

@app.route('/pdf_to_jpg', methods=['GET', 'POST'])
def pdf_to_jpg():
    if request.method == 'POST':
        # Check if a file was uploaded
        if 'file' not in request.files:
            return render_template('pdf_to_jpg.html', error='No file part')

        file = request.files['file']
        # Check if the file is empty
        if file.filename == '':
            return render_template('pdf_to_jpg.html', error='No selected file')

        # Check if the file is a PDF
        if file and allowed_file(file.filename):
            # Save the uploaded file
            filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filename)

            # Convert PDF to JPG images
            images = convert_from_path(filename)

            # Create a BytesIO object to hold the zip file in memory
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
                for i, image in enumerate(images):
                    # Save each image to a temporary file
                    temp_filename = os.path.join(app.config['UPLOAD_FOLDER'], f'page_{i+1}.jpg')
                    image.save(temp_filename, 'JPEG')
                    # Add the image file to the zip archive
                    zip_file.write(temp_filename, os.path.basename(temp_filename))
                    # Remove the temporary file
                    os.remove(temp_filename)

            # Move the buffer's position to the beginning to prepare for reading
            zip_buffer.seek(0)

            # Return the zip archive as a downloadable file
            return send_file(zip_buffer, attachment_filename='converted_images.zip', as_attachment=True)

    return render_template('pdf_to_jpg.html')


@app.route('/pdf_to_powerpoint', methods=['GET', 'POST'])
def pdf_to_powerpoint():
    if request.method == 'POST':
        if 'file' not in request.files:
            return render_template('pdf_to_powerpoint.html', error='No file part')

        file = request.files['file']
        if file.filename == '':
            return render_template('pdf_to_powerpoint.html', error='No selected file')

        if file and allowed_file(file.filename):
            # Save the uploaded file
            filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filename)

            # Convert PDF to PowerPoint
            images = convert_from_path(filename)
            pptx_filename = os.path.join(app.config['CONVERTED_FOLDER'], file.filename.replace('.pdf', '.pptx'))
            prs = Presentation()

            for i, image in enumerate(images):
                # Perform OCR on the image to extract text
                text = extract_text_from_image(image)

                # Add the text to a new slide
                slide_layout = prs.slide_layouts[5]  # Use a suitable slide layout
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f"Page {i+1}"
                content = slide.placeholders[1]
                content.text = text

            prs.save(pptx_filename)

            # Return the converted PowerPoint file for download
            return send_file(pptx_filename, as_attachment=True)

    return render_template('pdf_to_powerpoint.html')

def format_excel(ws, df):
    for r_idx, row in enumerate(df.iterrows(), start=1):
        for c_idx, value in enumerate(row[1], start=1):
            cell = ws.cell(row=r_idx, column=c_idx, value=value)
            cell.alignment = Alignment(horizontal='center', vertical='center')
            cell.font = Font(name='Arial', size=10)
            if r_idx == 1:  # Header row
                cell.fill = PatternFill(start_color="FFC000", end_color="FFC000", fill_type="solid")  # Yellow background for header
                cell.font = Font(bold=True)
            else:
                if c_idx % 2 == 0:  # Alternate row color
                    cell.fill = PatternFill(start_color="E6E6E6", end_color="E6E6E6", fill_type="solid")

@app.route('/pdf_to_excel', methods=['GET', 'POST'])
def pdf_to_excel():
    if request.method == 'POST':
        if 'file' not in request.files:
            return render_template('pdf_to_excel.html', error='No file part')

        file = request.files['file']
        if file.filename == '':
            return render_template('pdf_to_excel.html', error='No selected file')

        if file and allowed_file(file.filename):
            # Save the uploaded file
            filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filename)

            # Convert PDF to CSV
            csv_filename = os.path.join(app.config['CONVERTED_FOLDER'], file.filename.replace('.pdf', '.csv'))
            try:
                tabula.convert_into(filename, csv_filename, output_format="csv", pages='all')
            except Exception as e:
                return render_template('pdf_to_excel.html', error=f'Error converting PDF to CSV: {str(e)}')

            # Read the CSV data into a pandas DataFrame
            df = pd.read_csv(csv_filename)

            # Create a new Excel workbook and sheet
            excel_filename = os.path.join(app.config['CONVERTED_FOLDER'], file.filename.replace('.pdf', '.xlsx'))
            wb = Workbook()
            ws = wb.active
            ws.title = 'Data'

            # Format the Excel sheet
            format_excel(ws, df)

            # Save the Excel file
            wb.save(excel_filename)

            # Return the converted Excel file for download
            return send_file(excel_filename, as_attachment=True)

    return render_template('pdf_to_excel.html')

if __name__ == '__main__':
    app.run(debug=True)